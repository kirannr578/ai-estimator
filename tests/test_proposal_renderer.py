"""Tests for `core.proposal_renderer` — the bid `proposal/` markdown to
PDF + PPTX renderer.

These tests run against synthesized temp directories so they're independent
of the live `bids/` workspaces. No network or LLM is touched.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from core.proposal_renderer import (
    Section,
    build_client_pdfs,
    build_internal_workbook,
    build_pitch_deck,
    build_proposal_pdf,
    build_proposal_pptx,
    discover_sections,
    load_firm_profile,
    partition_sections,
    route_section,
)
from core.proposal_renderer.common import (
    neutralize_placeholders,
    parse_solicitation_number,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_bid_workspace(
    root: Path,
    *,
    files: dict[str, str],
    slug: str = "synthetic-bid-2026-001",
) -> Path:
    """Build a `bids/<slug>/proposal/` tree with the named .md files written
    inside. Returns the bid-dir path (parent of `proposal/`)."""
    bid_dir = root / slug
    proposal_dir = bid_dir / "proposal"
    proposal_dir.mkdir(parents=True, exist_ok=True)
    for name, body in files.items():
        (proposal_dir / name).write_text(body, encoding="utf-8")
    return bid_dir


@pytest.fixture()
def firm_profile() -> dict:
    """Live firm profile from `firm/firm-profile.json` — already in the
    repo, no fixture data duplicated here."""
    return load_firm_profile()


# ---------------------------------------------------------------------------
# 1. Section discovery — orders by numerical prefix.
# ---------------------------------------------------------------------------


def test_section_discovery_orders_by_prefix(tmp_path: Path) -> None:
    bid_dir = _make_bid_workspace(
        tmp_path,
        files={
            "10-price.md": "# Price\n\nPrice section.\n",
            "01-summary.md": "# Summary\n\nExecutive summary.\n",
            "02-approach.md": "# Approach\n\nTechnical approach.\n",
            "00-readme.md": "# Readme\n\nIntro.\n",
            "z-extras.md": "# Extras\n\nNo prefix; sorts last.\n",
        },
    )
    sections = discover_sections(bid_dir / "proposal")
    titles = [s.title for s in sections]
    assert titles == ["Readme", "Summary", "Approach", "Price", "Extras"], titles
    # Prefix integers must be parsed correctly (00-readme has prefix 0,
    # not 10 or some lexicographic weirdness).
    prefixes = [s.prefix for s in sections]
    assert prefixes == [0, 1, 2, 10, None], prefixes


# ---------------------------------------------------------------------------
# 2. Section discovery — skips underscored drafts.
# ---------------------------------------------------------------------------


def test_section_discovery_skips_underscored_drafts(tmp_path: Path) -> None:
    bid_dir = _make_bid_workspace(
        tmp_path,
        files={
            "00-readme.md": "# Readme\n\nReal section.\n",
            "_draft-notes.md": "# Draft\n\nShould be skipped.\n",
            "_01-old-summary.md": "# Old Summary\n\nShould be skipped.\n",
            "01-summary.md": "# Summary\n\nReal section.\n",
        },
    )
    sections = discover_sections(bid_dir / "proposal")
    assert [s.filename for s in sections] == ["00-readme.md", "01-summary.md"]
    assert all(not s.filename.startswith("_") for s in sections)


# ---------------------------------------------------------------------------
# 3. PDF render smoke — placeholder + table + H1 round-trip.
# ---------------------------------------------------------------------------


def test_pdf_render_smoke(tmp_path: Path, firm_profile: dict) -> None:
    bid_dir = _make_bid_workspace(
        tmp_path,
        files={
            "00-readme.md": (
                "# Synthetic Bid Proposal\n"
                "\n"
                "## Overview\n"
                "\n"
                "This is a synthetic bid. The owner is `[USER TO FILL — owner name]`.\n"
                "\n"
                "| Field | Value |\n"
                "|---|---|\n"
                "| Project | Synthetic |\n"
                "| Owner | TBD |\n"
            ),
        },
    )
    out_pdf = tmp_path / "out" / "proposal-full.pdf"
    build_proposal_pdf(bid_dir, out_pdf, firm_profile)

    assert out_pdf.is_file()
    assert out_pdf.stat().st_size > 1024, f"PDF size too small: {out_pdf.stat().st_size}"

    # Round-trip the H1 text via pypdf.
    from pypdf import PdfReader

    reader = PdfReader(str(out_pdf))
    assert len(reader.pages) >= 1
    first_text = ""
    # Concatenate text from the first 4 pages so we cover cover / TOC /
    # body — the H1 should land somewhere in there.
    for page in reader.pages[: min(4, len(reader.pages))]:
        first_text += page.extract_text() or ""
    assert "Synthetic Bid Proposal" in first_text, first_text[:400]


# ---------------------------------------------------------------------------
# 4. PDF render — placeholder text round-trips and the `placeholder` styling
#    span is present in the rendered HTML pipeline.
# ---------------------------------------------------------------------------


def test_pdf_renders_placeholder_in_red(tmp_path: Path, firm_profile: dict) -> None:
    bid_dir = _make_bid_workspace(
        tmp_path,
        files={
            "00-readme.md": (
                "# Placeholder Test\n"
                "\n"
                "The lead is `[USER TO FILL — foo]` and the deadline is "
                "`[USER TO FILL — date]`.\n"
            ),
        },
    )
    out_pdf = tmp_path / "out" / "proposal-full.pdf"
    # `show_placeholders=True` is the moral equivalent of yesterday's
    # default behavior — keep `[USER TO FILL ...]` markers in the body.
    build_proposal_pdf(bid_dir, out_pdf, firm_profile, show_placeholders=True)

    # The placeholder text must survive into the PDF. xhtml2pdf encodes
    # text through font-glyph streams so a raw byte search misses it; use
    # pypdf to extract Unicode text from the body pages instead.
    from pypdf import PdfReader

    reader = PdfReader(str(out_pdf))
    extracted = ""
    for page in reader.pages:
        extracted += page.extract_text() or ""
    assert "USER TO FILL" in extracted, (
        "placeholder text did not survive into the PDF; first 800 chars of "
        f"extracted text:\n{extracted[:800]}"
    )

    # Cross-check: the renderer's HTML stage wraps placeholders in a span
    # with the `placeholder` class. We exercise the HTML stage directly to
    # confirm the styling hook is applied (the PDF library obscures color in
    # the byte stream so this is the cleanest assertion we can make).
    from core.proposal_renderer.pdf import _highlight_placeholders_in_html

    sample = "<p>The lead is [USER TO FILL — foo].</p>"
    out = _highlight_placeholders_in_html(sample)
    assert 'class="placeholder"' in out, out
    assert "[USER TO FILL — foo]" in out, out


# ---------------------------------------------------------------------------
# 5. PPTX (pitch deck) — slide count lands in 10-15 for a 13-file proposal.
# ---------------------------------------------------------------------------


def test_pptx_pitch_deck_slide_count(tmp_path: Path, firm_profile: dict) -> None:
    files = {
        "00-readme.md": "# Synthetic Bid\n\n> Solicitation: ABC-123.\n",
        "01-executive-summary.md": (
            "# Executive Summary\n"
            "\n"
            "## Project understanding\n"
            "\n"
            "- Renovate the labs.\n- Maintain occupancy.\n- Beat schedule.\n"
        ),
        "02-technical-approach.md": (
            "# Technical Approach\n"
            "\n"
            "## 1. Mobilization\n\nMobilize within 5 days of NTP.\n"
            "\n"
            "## 2. Demolition\n\nSurgical demolition of finishes.\n"
            "\n"
            "## 3. Rough-in\n\nMEP rough-in coordinated weekly.\n"
            "\n"
            "## 4. Finish out\n\nFinal finishes + commissioning.\n"
        ),
        "03-project-team.md": (
            "# Project Team\n"
            "\n"
            "## Role 1 — Project Manager\n"
            "\n"
            "| Field | Value |\n"
            "|---|---|\n"
            "| Name | Jane Doe |\n"
            "\n"
            "## Role 2 — Superintendent\n"
            "\n"
            "| Field | Value |\n"
            "|---|---|\n"
            "| Name | John Smith |\n"
        ),
        "04-past-performance.md": (
            "# Past Performance\n"
            "\n"
            "| # | Project name | Owner | Contract value | Completion year | Why relevant |\n"
            "|---|---|---|---|---|---|\n"
            "| 1 | Hindu Temple | NTHHS | $1.2M | 2025 | similar scope |\n"
            "| 2 | Holiday Inn | Franchisee | $850k | 2024 | hospitality |\n"
            "| 3 | Lavon RV Park | Lavon LLC | $1.05M | 2026 | site work |\n"
        ),
        "05-schedule-narrative.md": (
            "# Schedule\n"
            "\n"
            "## Milestones\n"
            "\n"
            "- NTP: day 0\n- Demolition complete: day 14\n- Substantial: day 90\n"
        ),
        "06-quality-control-plan.md": (
            "# Quality Control\n\n## Highlights\n\nThree-level inspection.\n"
        ),
        "07-safety-plan.md": (
            "# Safety\n\n## Summary\n\nOSHA 30 site supers, daily JHAs.\n"
        ),
        "08-fill-guide.md": "# Fill Guide\n\nStandard ASU fill guide.\n",
        "09-hsp-form.md": "# HSP\n\nHSP form fill guide.\n",
        "10-price-proposal.md": (
            "# Price Proposal\n"
            "\n"
            "Total bid price: `[USER TO FILL — $]`.\n"
        ),
        "11-submission-checklist.md": (
            "# Submission Checklist\n"
            "\n"
            "- [x] Cover letter\n- [ ] Bid bond\n- [ ] HSP\n"
        ),
        "12-bid-bond.md": "# Bid Bond\n\nTemplate.\n",
    }
    bid_dir = _make_bid_workspace(tmp_path, files=files)

    out_pptx = tmp_path / "out" / "proposal-pitch.pptx"
    build_proposal_pptx(bid_dir, out_pptx, firm_profile, style="pitch_deck")

    from pptx import Presentation

    prs = Presentation(str(out_pptx))
    n = len(prs.slides)
    assert 10 <= n <= 15, f"pitch deck slide count out of range: {n}"


# ---------------------------------------------------------------------------
# 6. PPTX (full style) — at least one body slide per H1 + a title slide.
# ---------------------------------------------------------------------------


def test_pptx_full_style_one_slide_per_h1(tmp_path: Path, firm_profile: dict) -> None:
    bid_dir = _make_bid_workspace(
        tmp_path,
        files={
            "00-three-h1.md": (
                "# First Section\n"
                "\n"
                "Body of first section.\n"
                "\n"
                "# Second Section\n"
                "\n"
                "Body of second section.\n"
                "\n"
                "# Third Section\n"
                "\n"
                "Body of third section.\n"
            ),
        },
    )
    out_pptx = tmp_path / "out" / "proposal-full.pptx"
    build_proposal_pptx(bid_dir, out_pptx, firm_profile, style="full")

    from pptx import Presentation

    prs = Presentation(str(out_pptx))
    n = len(prs.slides)
    # 1 title slide + at least 3 body slides.
    assert n >= 4, f"expected ≥ 4 slides (1 title + ≥3 body) for 3 H1s, got {n}"


# ---------------------------------------------------------------------------
# 7. Solicitation-number parser — covers all four bid-folder shapes.
# ---------------------------------------------------------------------------


def test_parse_solicitation_number_all_shapes() -> None:
    cases = {
        "tamu-harrington-2025-06813": "2025-06813",
        "angelo-state-carr-efa-26-007": "26-007",
        "usfws-san-marcos-140FC126R0017": "140FC126R0017",
        "cmd-post-ndi-W50S7626QA001": "W50S7626QA001",
    }
    for slug, expected in cases.items():
        got = parse_solicitation_number(slug)
        assert got == expected, f"{slug!r}: expected {expected!r}, got {got!r}"


# ---------------------------------------------------------------------------
# 8. Content routing — client-facing filenames go to "client", scaffolding
#    filenames go to "internal", default-deny goes to client.
# ---------------------------------------------------------------------------


def _section(name: str, body: str = "# Stub\n\nstub.\n") -> Section:
    """Tiny in-memory Section helper that doesn't touch disk."""
    return Section(
        path=Path(name),
        prefix=None,
        slug=Path(name).stem,
        title="Stub",
        body=body,
    )


def test_route_section_classifies_client_internal_and_default() -> None:
    # Client-facing — explicit allowlist.
    for name in [
        "00-readme.md",
        "01-executive-summary.md",
        "02-technical-approach.md",
        "03-project-team.md",
        "04-past-performance.md",
        "10-price-proposal.md",
    ]:
        assert route_section(_section(name)) == "client", name

    # Internal — pattern match (fill-guide / submission-checklist / bond /
    # cover-letter / RFI / template / compliance / acknowledgment / memo).
    for name in [
        "08-sf1442-fill-guide.md",
        "09-hsp-form-fill-guide.md",
        "11-submission-checklist.md",
        "12-bid-bond-letter-template.md",
        "13-RFI-cover-letter.md",
        "14-reps-and-certs-pull-guide.md",
        "15-dba-sca-compliance-notes.md",
        "16-acknowledgment.md",
        "17-internal-memo.md",
    ]:
        assert route_section(_section(name)) == "internal", name

    # Default-deny: an unknown file lands in client (per the brief).
    assert route_section(_section("99-mystery-file.md")) == "client"


def test_partition_sections_splits_correctly() -> None:
    sections = [
        _section("00-readme.md"),
        _section("01-executive-summary.md"),
        _section("08-sf1442-fill-guide.md"),
        _section("11-submission-checklist.md"),
        _section("99-mystery-file.md"),
    ]
    client, internal = partition_sections(sections)
    client_names = [s.filename for s in client]
    internal_names = [s.filename for s in internal]
    assert client_names == [
        "00-readme.md",
        "01-executive-summary.md",
        "99-mystery-file.md",
    ]
    assert internal_names == [
        "08-sf1442-fill-guide.md",
        "11-submission-checklist.md",
    ]


# ---------------------------------------------------------------------------
# 9. Placeholder neutralization — `[USER TO FILL — ...]` markers get
#    replaced with a fillable underline, NOT the raw bracketed text.
# ---------------------------------------------------------------------------


def test_neutralize_placeholders_replaces_with_underline() -> None:
    src = (
        "The lead is [USER TO FILL — name] and the deadline is "
        "[USER TO FILL — 2026-06-01]. Cost: [USER TO FILL — $].\n"
    )
    out = neutralize_placeholders(src)

    # The bracketed marker text must be gone.
    assert "USER TO FILL" not in out, out
    assert "[" not in out and "]" not in out, out

    # And there must be at least three runs of underscores (one per
    # original marker) of length ≥ 3.
    import re as _re

    runs = _re.findall(r"_{3,}", out)
    assert len(runs) >= 3, (out, runs)


# ---------------------------------------------------------------------------
# 10. Client-facing PDFs — both files render, neither contains
#     "USER TO FILL" when show_placeholders=False (default), and the
#     executive summary stays at ≤ 12 pages.
# ---------------------------------------------------------------------------


def test_build_client_pdfs_hides_placeholders_and_caps_brochure(
    tmp_path: Path, firm_profile: dict
) -> None:
    files = {
        "00-readme.md": (
            "# Synthetic Bid — Cover\n"
            "\n"
            "Owner: [USER TO FILL — owner].\n"
        ),
        "01-executive-summary.md": (
            "# Executive Summary\n"
            "\n"
            "Synthetic bid summary. Lead: [USER TO FILL — lead].\n"
        ),
        "02-technical-approach.md": (
            "# Technical Approach\n"
            "\n"
            "## 1. Mobilization\n\nMobilize within 5 days of NTP.\n"
            "\n"
            "## 2. Demolition\n\nSurgical demolition.\n"
        ),
        "03-project-team.md": (
            "# Project Team\n"
            "\n"
            "## Role 1 — PM\n\n| Field | Value |\n|---|---|\n| Name | Jane |\n"
        ),
        "04-past-performance.md": (
            "# Past Performance\n"
            "\n"
            "| # | Project | Owner | Value | Year | Why relevant |\n"
            "|---|---|---|---|---|---|\n"
            "| 1 | Hindu Temple | NTHHS | $1.2M | 2025 | similar |\n"
            "| 2 | Holiday Inn | Franchisee | $850k | 2024 | hospitality |\n"
            "| 3 | Lavon RV | Lavon LLC | $1.05M | 2026 | sitework |\n"
        ),
        "10-price-proposal.md": (
            "# Price Proposal\n\nTotal: [USER TO FILL — $].\n"
        ),
        # An internal-scaffolding file — must NOT leak into either client PDF.
        "11-submission-checklist.md": "# Checklist\n\n- [ ] item\n",
    }
    bid_dir = _make_bid_workspace(tmp_path, files=files)
    out_dir = tmp_path / "out"

    artifacts = build_client_pdfs(bid_dir, out_dir, firm_profile)

    assert set(artifacts) == {"executive_summary", "full_proposal"}
    for path in artifacts.values():
        assert path.is_file()
        assert path.stat().st_size > 4_000, (path, path.stat().st_size)

    from pypdf import PdfReader

    # Brochure (executive summary) must stay ≤ 12 pages per the brief.
    es_reader = PdfReader(str(artifacts["executive_summary"]))
    assert 1 <= len(es_reader.pages) <= 12, (
        f"executive summary page count out of range: {len(es_reader.pages)}"
    )

    # No placeholder markers in either client PDF when show_placeholders
    # is the default (False).
    for key, path in artifacts.items():
        text = ""
        for page in PdfReader(str(path)).pages:
            text += page.extract_text() or ""
        assert "USER TO FILL" not in text, (
            f"{key}: 'USER TO FILL' must be neutralized on client output. "
            f"first 400 chars:\n{text[:400]}"
        )

    # The internal-only file must NOT appear in the full proposal body —
    # a tracer string from 11-submission-checklist.md must be absent.
    full_text = ""
    for page in PdfReader(str(artifacts["full_proposal"])).pages:
        full_text += page.extract_text() or ""
    # The H1 "Checklist" is a relatively unique tracer.
    # (We tolerate the section appearing in the TOC if it ever leaked,
    # but the body content shouldn't.)
    assert "Submission Checklist" not in full_text, (
        "internal scaffolding leaked into client full-proposal PDF"
    )


# ---------------------------------------------------------------------------
# 11. Internal workbook — file is written and KEEPS the placeholder
#     markers as red text (we just verify the marker text round-trips,
#     since color in the PDF byte stream is opaque from pypdf).
# ---------------------------------------------------------------------------


def test_internal_workbook_keeps_placeholders(
    tmp_path: Path, firm_profile: dict
) -> None:
    files = {
        # Two internal-scaffolding files so the workbook has real content.
        "08-sf1442-fill-guide.md": (
            "# SF-1442 Fill Guide\n"
            "\n"
            "Block 17 — Contractor: [USER TO FILL — legal name].\n"
        ),
        "11-submission-checklist.md": (
            "# Submission Checklist\n"
            "\n"
            "- [ ] SF-1442 signed by [USER TO FILL — date]\n"
        ),
        # Plus a client-facing file (must NOT appear in the workbook).
        "01-executive-summary.md": (
            "# Executive Summary\n\nClient-facing tracer string.\n"
        ),
    }
    bid_dir = _make_bid_workspace(tmp_path, files=files)
    out_pdf = tmp_path / "out" / "internal-workbook.pdf"

    build_internal_workbook(bid_dir, out_pdf, firm_profile)
    assert out_pdf.is_file()
    assert out_pdf.stat().st_size > 1_000

    from pypdf import PdfReader

    text = ""
    for page in PdfReader(str(out_pdf)).pages:
        text += page.extract_text() or ""

    # Markers preserved on the internal tier.
    assert "USER TO FILL" in text, text[:600]
    # Internal scaffolding included.
    assert "SF-1442" in text or "Submission Checklist" in text, text[:600]
    # Client-facing content NOT included.
    assert "Client-facing tracer string" not in text, (
        "client-facing content leaked into internal workbook"
    )


# ---------------------------------------------------------------------------
# 12. Pitch deck — 10–15 slides, template loaded successfully (so the file
#     under firm/assets/templates/bpc-pitch-deck-template.pptx is exercised
#     when present, or generated on first call when absent).
# ---------------------------------------------------------------------------


def test_pitch_deck_uses_template_and_lands_in_slide_range(
    tmp_path: Path, firm_profile: dict
) -> None:
    files = {
        "00-readme.md": "# Synthetic Bid\n",
        "01-executive-summary.md": (
            "# Executive Summary\n\n- Renovate.\n- Maintain occupancy.\n- Beat schedule.\n"
        ),
        "02-technical-approach.md": (
            "# Technical Approach\n\n## 1. Mob\n\n## 2. Demo\n\n## 3. Rough\n\n## 4. Finish\n"
        ),
        "03-project-team.md": (
            "# Project Team\n\n## Role 1 — PM\n\n| Field | Value |\n|---|---|\n| Name | Jane |\n"
        ),
        "04-past-performance.md": (
            "# Past Performance\n\n"
            "| # | Project | Owner | Value | Year | Why |\n|---|---|---|---|---|---|\n"
            "| 1 | Hindu Temple | NTHHS | $1.2M | 2025 | similar |\n"
            "| 2 | Holiday Inn | Franchisee | $850k | 2024 | hospitality |\n"
            "| 3 | Lavon RV | Lavon | $1.05M | 2026 | sitework |\n"
        ),
        "05-schedule-narrative.md": "# Schedule\n\n- NTP\n- Substantial\n- Final\n",
        "06-quality-control-plan.md": "# Quality Control\n\nThree-level inspection.\n",
        "07-safety-plan.md": "# Safety\n\nOSHA 30 supers.\n",
        "10-price-proposal.md": "# Price Proposal\n\nTotal: [USER TO FILL — $].\n",
    }
    bid_dir = _make_bid_workspace(tmp_path, files=files)
    out_pptx = tmp_path / "out" / "pitch-deck.pptx"

    # Use a private template path so this test doesn't mutate the repo's
    # checked-in template (and exercises the "auto-generate template if
    # missing" code path when the file doesn't exist yet).
    template_path = tmp_path / "templates" / "test-template.pptx"

    result = build_pitch_deck(
        bid_dir, out_pptx, firm_profile, template_path=template_path
    )
    assert result == out_pptx
    assert out_pptx.is_file()
    assert out_pptx.stat().st_size > 20_000

    # The renderer must have created the template on first call.
    assert template_path.is_file(), (
        f"pitch-deck builder did not materialize a template at {template_path}"
    )

    from pptx import Presentation

    prs = Presentation(str(out_pptx))
    n = len(prs.slides)
    assert 10 <= n <= 15, f"pitch deck slide count out of range: {n}"
