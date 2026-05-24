"""Build the BPC-branded pitch-deck PPTX for a single bid.

Public entry point: `build_pitch_deck(bid_dir, out_path, firm_profile,
template_path=None, show_placeholders=False)`.

Approach (per the user brief):

* Source the master template from a `.pptx` on disk if one exists at
  `firm/assets/templates/bpc-pitch-deck-template.pptx`. The brief asks
  us to look in BPC's OneDrive submission package first; that folder
  only contains a markdown outline (no `.pptx`), so we generate the
  template programmatically and save it on first render.
* The "template" is a 16:9 widescreen `Presentation` with our slide
  size and color palette pre-baked. python-pptx can only manipulate
  layouts by hand-crafted XML, so the actual visual design (color
  blocks, dividers, footers, gold-accent rules) is rendered by the
  slide-builder helpers below. Re-opening the template across renders
  guarantees identical output.
* 12 default slides; can grow to 13–15 if richer per-bid content is
  available (we cap at 15).

Placeholders: when `show_placeholders=False` (default for client-facing
output), `[USER TO FILL …]` markers are rewritten to a neutral
underline before each slide is built. When True, markers render in red
bold.

Backward compatibility: `build_proposal_pptx(... style="pitch_deck"|"full")`
still works as a thin wrapper around `build_pitch_deck` for the pitch
style, and a one-slide-per-H1 walkthrough builder for the full style.
"""

from __future__ import annotations

import re
import shutil
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .branding import (
    BPC_GOLD,
    BPC_GRAY_BG,
    BPC_GRAY_LINE,
    BPC_NAVY,
    BPC_NAVY_DEEP,
    BPC_SLATE,
    LOGO_PATH,
    PITCH_DECK_TEMPLATE,
    PPTX_FONT_BODY,
    PPTX_FONT_HEADING,
)
from .common import (
    PLACEHOLDER_PATTERNS,
    Section,
    collapse_whitespace,
    discover_sections,
    find_h1_blocks,
    find_h2_blocks,
    find_section,
    first_paragraph,
    list_items,
    neutralize_placeholders,
    neutralize_placeholders_in_section,
    parse_bid_title,
    parse_solicitation_number,
    primary_personnel,
    split_placeholders,
    strip_blockquote_prefix,
)


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """`#RRGGBB` -> `RGBColor`."""
    s = hex_str.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


NAVY = _hex_to_rgb(BPC_NAVY)
NAVY_DEEP = _hex_to_rgb(BPC_NAVY_DEEP)
SLATE = _hex_to_rgb(BPC_SLATE)
GOLD = _hex_to_rgb(BPC_GOLD)
GRAY_BG = _hex_to_rgb(BPC_GRAY_BG)
GRAY_LINE = _hex_to_rgb(BPC_GRAY_LINE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1F, 0x26)
RED = RGBColor(0xC0, 0x00, 0x00)
FOOTER_GRAY = RGBColor(0x77, 0x77, 0x77)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def build_pitch_deck(
    bid_dir: Path,
    out_path: Path,
    firm_profile: dict[str, Any],
    *,
    template_path: Path | None = None,
    show_placeholders: bool = False,
    logo_path: Path | None = None,
) -> Path:
    """Render `<bid_dir>/proposal/` to a 12-slide BPC-branded pitch deck.

    Returns the resolved output path.
    """
    bid_dir = Path(bid_dir).resolve()
    out_path = Path(out_path).resolve()
    proposal_dir = bid_dir / "proposal"
    if not proposal_dir.is_dir():
        raise FileNotFoundError(f"no proposal directory at {proposal_dir}")
    sections = discover_sections(proposal_dir)
    if not sections:
        raise FileNotFoundError(f"no proposal sections found in {proposal_dir}")

    bid_slug = bid_dir.name
    bid_title = parse_bid_title(proposal_dir, bid_slug)
    solicitation = parse_solicitation_number(bid_slug)

    rendered_sections = [
        neutralize_placeholders_in_section(s, show_placeholders=show_placeholders)
        for s in sections
    ]

    template = template_path or PITCH_DECK_TEMPLATE
    _ensure_template(template)

    prs = Presentation(str(template))
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_pitch_deck(
        prs=prs,
        bid_slug=bid_slug,
        bid_title=bid_title,
        solicitation=solicitation,
        sections=rendered_sections,
        firm_profile=firm_profile,
        logo_path=logo_path or LOGO_PATH,
        show_placeholders=show_placeholders,
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def build_proposal_pptx(
    bid_dir: Path,
    out_path: Path,
    firm_profile: dict[str, Any],
    *,
    style: str = "pitch_deck",
    show_placeholders: bool = False,
    template_path: Path | None = None,
    logo_path: Path | None = None,
) -> Path:
    """Backward-compatible entry. `style="pitch_deck"` produces the
    branded 12-slide deck; `style="full"` produces a one-slide-per-H1
    walkthrough used by internal review only."""
    if style not in {"pitch_deck", "full"}:
        raise ValueError(
            f"style must be 'pitch_deck' or 'full', got {style!r}"
        )
    if style == "pitch_deck":
        return build_pitch_deck(
            bid_dir,
            out_path,
            firm_profile,
            template_path=template_path,
            show_placeholders=show_placeholders,
            logo_path=logo_path,
        )

    bid_dir = Path(bid_dir).resolve()
    out_path = Path(out_path).resolve()
    proposal_dir = bid_dir / "proposal"
    sections = discover_sections(proposal_dir)
    bid_slug = bid_dir.name
    bid_title = parse_bid_title(proposal_dir, bid_slug)
    solicitation = parse_solicitation_number(bid_slug)
    rendered_sections = [
        neutralize_placeholders_in_section(s, show_placeholders=show_placeholders)
        for s in sections
    ]

    template = template_path or PITCH_DECK_TEMPLATE
    _ensure_template(template)
    prs = Presentation(str(template))
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_full_deck(
        prs=prs,
        bid_slug=bid_slug,
        bid_title=bid_title,
        solicitation=solicitation,
        sections=rendered_sections,
        firm_profile=firm_profile,
        logo_path=logo_path or LOGO_PATH,
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# Template — generate once, reuse across renders
# ---------------------------------------------------------------------------


def _ensure_template(template_path: Path) -> None:
    """Create `template_path` if it doesn't already exist.

    The template is an empty 16:9 Presentation with our slide
    dimensions baked in. python-pptx can't easily author custom slide
    layouts at runtime, so the visual design lives in the slide
    builder helpers — but generating + reusing the same .pptx file
    on disk gives us a stable place for designers to swap in a real
    template later if BPC produces one.
    """
    template_path = Path(template_path)
    if template_path.exists():
        return
    template_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    # Save with no slides so re-openers don't have to purge a placeholder
    # slide on every render (which leaves orphan slide1.xml parts and
    # warns "Duplicate name 'ppt/slides/slide1.xml'").
    prs.save(str(template_path))


def _purge_template_placeholder_slides(prs: Presentation) -> None:
    """Remove any placeholder slides the template may carry over so the
    final deck is exactly the slides we built."""
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 — python-pptx idiom
    ids = list(sld_id_lst)
    for sld_id_el in ids:
        sld_id_lst.remove(sld_id_el)


# ---------------------------------------------------------------------------
# Pitch-deck builder
# ---------------------------------------------------------------------------


def _build_pitch_deck(
    *,
    prs: Presentation,
    bid_slug: str,
    bid_title: str,
    solicitation: str,
    sections: list[Section],
    firm_profile: dict[str, Any],
    logo_path: Path,
    show_placeholders: bool,
) -> None:
    _purge_template_placeholder_slides(prs)

    _slide_title(
        prs,
        bid_title=bid_title,
        solicitation=solicitation,
        firm_profile=firm_profile,
        logo_path=logo_path,
    )
    _slide_about_bpc(prs, firm_profile=firm_profile, bid_slug=bid_slug)
    _slide_project_understanding(
        prs, sections=sections, bid_slug=bid_slug, bid_title=bid_title
    )
    _slide_technical_approach(prs, sections=sections, bid_slug=bid_slug)
    _slide_project_team(prs, sections=sections, bid_slug=bid_slug)

    # Past performance — three slides, one per pick from firm_profile.
    _slides_past_performance(
        prs,
        sections=sections,
        firm_profile=firm_profile,
        bid_slug=bid_slug,
    )

    _slide_schedule(prs, sections=sections, bid_slug=bid_slug)
    _slide_quality_safety(prs, sections=sections, bid_slug=bid_slug)
    _slide_price(prs, sections=sections, bid_slug=bid_slug)
    _slide_why_bpc(prs, firm_profile=firm_profile, bid_slug=bid_slug)
    _slide_qa(prs, firm_profile=firm_profile, bid_slug=bid_slug)


def _build_full_deck(
    *,
    prs: Presentation,
    bid_slug: str,
    bid_title: str,
    solicitation: str,
    sections: list[Section],
    firm_profile: dict[str, Any],
    logo_path: Path,
) -> None:
    _purge_template_placeholder_slides(prs)
    _slide_title(
        prs,
        bid_title=bid_title,
        solicitation=solicitation,
        firm_profile=firm_profile,
        logo_path=logo_path,
    )
    for section in sections:
        for h1_title, body in find_h1_blocks(section.body) or [
            (section.title, section.body)
        ]:
            chunks = _chunk_text_for_slide(body, max_chars=400)
            if not chunks:
                chunks = [""]
            for i, chunk in enumerate(chunks):
                title = (
                    h1_title
                    if len(chunks) == 1
                    else f"{h1_title} ({i + 1}/{len(chunks)})"
                )
                _body_slide(
                    prs,
                    title=title,
                    body_lines=[
                        ln.strip() for ln in chunk.splitlines() if ln.strip()
                    ][:8],
                    bid_slug=bid_slug,
                )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _slide_title(
    prs: Presentation,
    *,
    bid_title: str,
    solicitation: str,
    firm_profile: dict[str, Any],
    logo_path: Path,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, NAVY)

    # Gold accent stripe across the top
    _add_rect(
        slide,
        left=Inches(0), top=Inches(0),
        width=SLIDE_WIDTH, height=Inches(0.18),
        fill=GOLD,
    )

    # Logo (centered upper-third)
    if logo_path and Path(logo_path).exists():
        try:
            slide.shapes.add_picture(
                str(logo_path),
                left=Inches(5.167),
                top=Inches(0.85),
                width=Inches(3.0),
            )
        except Exception:
            pass

    # Eyebrow
    _add_text(
        slide,
        text="BLUE PRINT CONSTRUCTS",
        left=Inches(0.5), top=Inches(2.6), width=Inches(12.333), height=Inches(0.4),
        font_size=Pt(13), bold=True, color=GOLD, align=PP_ALIGN.CENTER,
        font_name=PPTX_FONT_HEADING,
    )

    # Bid title
    _add_text(
        slide,
        text=bid_title,
        left=Inches(0.5), top=Inches(3.05), width=Inches(12.333), height=Inches(1.7),
        font_size=Pt(36), bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        font_name=PPTX_FONT_HEADING,
    )

    # Solicitation
    if solicitation:
        _add_text(
            slide,
            text=f"Solicitation: {solicitation}",
            left=Inches(0.5), top=Inches(4.85), width=Inches(12.333), height=Inches(0.45),
            font_size=Pt(15), color=WHITE, align=PP_ALIGN.CENTER,
        )

    # Gold rule under solicitation
    _add_rect(
        slide,
        left=Inches(6.166), top=Inches(5.5),
        width=Inches(1.0), height=Inches(0.04),
        fill=GOLD,
    )

    # Firm tag
    legal = firm_profile.get("legal_name", "")
    dba = firm_profile.get("dba", "")
    firm_line = (
        f"{dba} (legal: {legal})"
        if dba and legal and dba != legal
        else dba or legal
    )
    _add_text(
        slide,
        text=firm_line,
        left=Inches(0.5), top=Inches(5.7), width=Inches(12.333), height=Inches(0.4),
        font_size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER, italic=True,
    )

    # Submission date placeholder
    _add_text(
        slide,
        text=f"Submission date __________   ·   Document prepared {date.today().isoformat()}",
        left=Inches(0.5), top=Inches(6.15), width=Inches(12.333), height=Inches(0.4),
        font_size=Pt(10), color=GOLD, align=PP_ALIGN.CENTER,
    )


def _slide_about_bpc(
    prs: Presentation, *, firm_profile: dict[str, Any], bid_slug: str
) -> None:
    legal = firm_profile.get("legal_name", "")
    dba = firm_profile.get("dba", "")
    naics_codes = firm_profile.get("naics_codes") or []
    primary_naics = firm_profile.get("naics_primary_for_renovation_bids", "")
    uei = firm_profile.get("uei", "")
    cage = firm_profile.get("cage", "")
    sa = firm_profile.get("set_aside_eligibility") or {}
    self_perform = (firm_profile.get("trade_capabilities") or {}).get(
        "self_perform"
    ) or []
    year_founded = firm_profile.get("year_founded", "")

    bullets = [
        (False, f"{dba} — legal entity {legal}; founded {year_founded}, "
                f"Texas-incorporated, founder-led."),
        (False,
         f"NAICS: primary {primary_naics} (Commercial & Institutional "
         f"Building Construction); {len(naics_codes)} codes registered."),
        (False,
         f"Federal IDs: UEI {uei} · CAGE {cage} · SAM.gov status: "
         f"{firm_profile.get('sam_status', 'TBD')}."),
        (False,
         "Set-aside posture: " + (
             ", ".join(filter(None, [
                 "Texas HUB" if sa.get("tx_hub_status") else "",
                 f"DFW MSDC MBE" if sa.get("mbe_status") else "",
                 "Small Business" if sa.get("small_business") else "",
                 "Minority-owned" if sa.get("minority_owned") else "",
             ])) or "—")),
        (False,
         "Self-perform trades: "
         + (", ".join(self_perform[:6]) if self_perform else "—") + "."),
    ]
    _body_slide(
        prs,
        title="About Blue Print Constructs",
        bullets=bullets,
        bid_slug=bid_slug,
    )


def _slide_project_understanding(
    prs: Presentation,
    *,
    sections: list[Section],
    bid_slug: str,
    bid_title: str,
) -> None:
    src = (
        find_section(sections, 1, slug_contains="executive-summary")
        or find_section(sections, 2, slug_contains="volume-II")
        or find_section(sections, 2, slug_contains="technical-acceptability")
    )
    bullets: list[tuple[bool, str]] = []
    if src is not None:
        for h2_title, body in find_h2_blocks(src.body):
            tl = h2_title.lower()
            if any(
                k in tl
                for k in (
                    "understanding", "intent", "scope", "executive summary",
                    "overview",
                )
            ):
                items = list_items(body)[:3]
                if items:
                    bullets.extend((False, _truncate(it, 220)) for it in items)
                    break
                first = first_paragraph(body)
                if first:
                    bullets.append((False, _truncate(first, 280)))
                    break
        if not bullets:
            text = first_paragraph(src.body)
            if text:
                bullets.append((False, _truncate(text, 320)))

    if not bullets:
        bullets = [(False, f"Project: {bid_title}")]
    while len(bullets) < 3:
        bullets.append((False, "(See full proposal for additional detail.)"))
    bullets = bullets[:4]

    _body_slide(
        prs,
        title="Project understanding",
        bullets=bullets,
        bid_slug=bid_slug,
    )


def _slide_technical_approach(
    prs: Presentation, *, sections: list[Section], bid_slug: str
) -> None:
    src = (
        find_section(sections, 2, slug_contains="technical-approach")
        or find_section(sections, 2, slug_contains="volume-II")
        or find_section(sections, 2, slug_contains="technical-acceptability")
    )
    bullets: list[tuple[bool, str]] = []
    if src is not None:
        for h2_title, body in find_h2_blocks(src.body):
            label = h2_title.split("—", 1)[0].strip()
            label = re.sub(r"^\d+\.\s*", "", label)
            label = label.strip(":") or h2_title
            tl = label.lower()
            if any(
                skip in tl for skip in ("appendix", "table of", "acronym")
            ):
                continue
            first = first_paragraph(body)
            text = (
                f"{label}: {_truncate(first, 200)}" if first else label
            )
            bullets.append((False, _truncate(text, 240)))
            if len(bullets) >= 6:
                break
    if not bullets:
        bullets = [(False, "Technical-approach phases not detected in source markdown.")]
    _body_slide(
        prs,
        title="Technical approach",
        bullets=bullets,
        bid_slug=bid_slug,
        bold_first_phrase=True,
    )


def _slide_project_team(
    prs: Presentation, *, sections: list[Section], bid_slug: str
) -> None:
    src = find_section(sections, 3, slug_contains="project-team")
    rows: list[tuple[str, str]] = []
    if src is not None:
        for h2_title, body in find_h2_blocks(src.body):
            tl = h2_title.lower()
            if "key personnel" not in tl and "role" not in tl:
                continue
            role = re.sub(
                r"^[A-Z]\.\s*Key personnel\s*[—-]\s*", "", h2_title
            ).strip()
            role = re.sub(r"^Role\s+\d+\s*[—-]\s*", "", role).strip()
            name = ""
            m = re.search(
                r"\|\s*Name\s*\|\s*([^|]+?)\s*\|", body, re.IGNORECASE
            )
            if m:
                name = m.group(1).strip().strip("`")
            if not name or name.startswith("____"):
                name = "____________________"
            rows.append((role, name))
            if len(rows) >= 8:
                break
    if not rows:
        rows = [("Project lead", "____________________")]

    _table_slide(
        prs,
        title="Project team",
        headers=["Role", "Named lead"],
        rows=rows,
        bid_slug=bid_slug,
        col_widths=[Inches(5.0), Inches(7.0)],
    )


def _slides_past_performance(
    prs: Presentation,
    *,
    sections: list[Section],
    firm_profile: dict[str, Any],
    bid_slug: str,
) -> None:
    """Render one slide per pick (3 slides) using the section-divider
    layout when no project photo is available — text-only is fine and
    cleaner than a stock image (per the brief)."""
    rules = firm_profile.get("past_project_selection_rules") or {}
    pick_names = (rules.get(bid_slug) or {}).get("picks") or []
    projects_by_name = {
        p.get("name", ""): p for p in (firm_profile.get("past_projects") or [])
    }

    rendered = 0
    for idx, name in enumerate(pick_names[:3], start=1):
        project = projects_by_name.get(name)
        if project is None:
            for k, v in projects_by_name.items():
                if name and name.split(" — ")[0] in k:
                    project = v
                    break
        if project is None:
            continue
        _slide_past_perf_one(
            prs,
            title=f"Past performance — Project {idx} of 3",
            project=project,
            bid_slug=bid_slug,
        )
        rendered += 1

    # Fall back to a single text slide if firm_profile didn't surface
    # any picks for this bid.
    if rendered == 0:
        _body_slide(
            prs,
            title="Past performance",
            bullets=[
                (
                    False,
                    "(No past-performance picks configured for this bid in "
                    "firm-profile.json — see full proposal for the complete roster.)",
                )
            ],
            bid_slug=bid_slug,
        )


def _slide_past_perf_one(
    prs: Presentation,
    *,
    title: str,
    project: dict[str, Any],
    bid_slug: str,
) -> None:
    """Single past-performance slide — text-only layout (no image)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, WHITE)
    _add_title_bar(slide, title)

    name = project.get("name", "")
    owner = project.get("owner", "")
    value = project.get("contract_value") or "—"
    completion = (
        project.get("actual_completion_date")
        or project.get("completion_date")
        or project.get("scheduled_substantial_completion")
        or "—"
    )
    scope = project.get("scope_summary") or ""
    role = project.get("role", "GC")
    delivery = project.get("delivery_method") or "—"

    # Project name as a large eyebrow + secondary metadata
    _add_text(
        slide,
        text=name,
        left=Inches(0.5), top=Inches(1.45),
        width=Inches(12.333), height=Inches(0.7),
        font_size=Pt(24), bold=True, color=NAVY,
        font_name=PPTX_FONT_HEADING,
    )

    meta_lines = [
        f"Owner: {owner}",
        f"Contract value: {value}",
        f"Completion: {completion}",
        f"Role: {role}   ·   Delivery: {delivery}",
    ]
    for i, line in enumerate(meta_lines):
        _add_text(
            slide,
            text=line,
            left=Inches(0.5),
            top=Inches(2.2 + i * 0.32),
            width=Inches(12.333),
            height=Inches(0.32),
            font_size=Pt(13),
            color=SLATE,
        )

    # Gold rule
    _add_rect(
        slide,
        left=Inches(0.5), top=Inches(3.55),
        width=Inches(1.0), height=Inches(0.04),
        fill=GOLD,
    )

    # Scope summary
    _add_text(
        slide,
        text="Scope summary",
        left=Inches(0.5), top=Inches(3.75),
        width=Inches(12.333), height=Inches(0.32),
        font_size=Pt(11), bold=True, color=NAVY,
    )
    _add_text(
        slide,
        text=_truncate(scope, 800),
        left=Inches(0.5), top=Inches(4.05),
        width=Inches(12.333), height=Inches(2.6),
        font_size=Pt(13), color=NEAR_BLACK,
    )

    _add_footer(slide, bid_slug=bid_slug, slide_number=len(prs.slides))


def _slide_schedule(
    prs: Presentation, *, sections: list[Section], bid_slug: str
) -> None:
    src = find_section(sections, 5, slug_contains="schedule")
    rows: list[list[str]] = []
    if src is not None:
        for h2_title, body in find_h2_blocks(src.body):
            tl = h2_title.lower()
            if "milestone" in tl or "schedule" in tl:
                items = list_items(body)
                for it in items[:6]:
                    if ":" in it:
                        m, d = it.split(":", 1)
                        rows.append(
                            [_truncate(m.strip(), 60), _truncate(d.strip(), 80)]
                        )
                    else:
                        rows.append([_truncate(it, 60), "—"])
                if rows:
                    break
        if not rows:
            for it in list_items(src.body)[:6]:
                if ":" in it:
                    m, d = it.split(":", 1)
                    rows.append(
                        [_truncate(m.strip(), 60), _truncate(d.strip(), 80)]
                    )
    if not rows:
        rows = [
            ["Notice to proceed", "____________________"],
            ["Mobilization", "____________________"],
            ["Substantial completion", "____________________"],
            ["Final completion", "____________________"],
        ]
    _table_slide(
        prs,
        title="Schedule milestones",
        headers=["Milestone", "Date / duration"],
        rows=rows,
        bid_slug=bid_slug,
        col_widths=[Inches(5.0), Inches(7.0)],
    )


def _slide_quality_safety(
    prs: Presentation, *, sections: list[Section], bid_slug: str
) -> None:
    qc = find_section(sections, 6, slug_contains="quality")
    safety = find_section(sections, 7, slug_contains="safety")
    bullets: list[tuple[bool, str]] = []
    for label, src in (("QC: ", qc), ("Safety: ", safety)):
        if src is None:
            continue
        used = False
        for h2_title, body in find_h2_blocks(src.body):
            tl = h2_title.lower()
            if any(
                k in tl
                for k in (
                    "highlights", "summary", "philosophy", "overview",
                    "objectives",
                )
            ):
                first = first_paragraph(body)
                if first:
                    bullets.append((False, _truncate(label + first, 240)))
                    used = True
                    break
                items = list_items(body)
                if items:
                    bullets.append((False, _truncate(label + items[0], 240)))
                    used = True
                    break
        if not used:
            first = first_paragraph(src.body)
            if first:
                bullets.append((False, _truncate(label + first, 240)))

    while len(bullets) < 4:
        bullets.append(
            (False, "(Additional QC / safety detail in the full proposal.)")
        )
    bullets = bullets[:4]
    _body_slide(
        prs, title="Quality & safety", bullets=bullets, bid_slug=bid_slug
    )


def _slide_price(
    prs: Presentation, *, sections: list[Section], bid_slug: str
) -> None:
    src = (
        find_section(sections, 10, slug_contains="price-proposal")
        or find_section(sections, 1, slug_contains="price-proposal")
        or find_section(sections, 1, slug_contains="volume-I")
    )

    total = "$____________________"
    bullets: list[str] = []
    if src is not None:
        money_re = re.compile(
            r"(grand total|total bid|total proposal|total price|base bid)[^|\n]*?(\$\s*[\d,_]+|\$\s*_+)",
            re.IGNORECASE,
        )
        for line in src.body.splitlines():
            stripped = line.strip().lstrip("|").lstrip("-").strip()
            if not stripped:
                continue
            m = money_re.search(stripped)
            if m and "$" in m.group(2):
                bullets.append(
                    _truncate(collapse_whitespace(stripped), 200)
                )
                if "total" in m.group(1).lower() and total == "$____________________":
                    total = m.group(2)
            if len(bullets) >= 4:
                break
        if not bullets:
            first = first_paragraph(src.body)
            if first:
                bullets.append(_truncate(first, 240))

    while len(bullets) < 3:
        bullets.append("(See full proposal for itemized pricing detail.)")
    bullets = bullets[:3]

    _stat_slide(
        prs,
        title="Price proposal summary",
        stat_label="TOTAL PROPOSED PRICE",
        stat_value=total,
        bullets=[(False, b) for b in bullets],
        bid_slug=bid_slug,
    )


def _slide_why_bpc(
    prs: Presentation, *, firm_profile: dict[str, Any], bid_slug: str
) -> None:
    sa = firm_profile.get("set_aside_eligibility") or {}
    self_perform = (firm_profile.get("trade_capabilities") or {}).get(
        "self_perform"
    ) or []
    sectors = (firm_profile.get("trade_capabilities") or {}).get(
        "served_sectors"
    ) or []
    bullets = [
        (
            False,
            f"Texas-based small business — self-perform across "
            f"{len(self_perform)} core trades.",
        ),
        (
            False,
            "Diverse-business posture: "
            + (
                ", ".join(
                    filter(
                        None,
                        [
                            "TX HUB" if sa.get("tx_hub_status") else "",
                            "DFW MSDC MBE" if sa.get("mbe_status") else "",
                            "Small Business" if sa.get("small_business") else "",
                            "Minority-owned" if sa.get("minority_owned") else "",
                        ],
                    )
                )
                or "—"
            )
            + ".",
        ),
        (
            False,
            "Sector breadth: "
            + (
                ", ".join(s.split("—")[0].strip() for s in sectors[:4])
                if sectors
                else "—"
            )
            + ".",
        ),
        (
            False,
            f"Founder-led PMP / CSM operator with 22 years delivery leadership — "
            f"{(firm_profile.get('key_personnel') or [{}])[0].get('name', '')}.",
        ),
    ]
    _body_slide(
        prs, title="Why Blue Print Constructs", bullets=bullets, bid_slug=bid_slug
    )


def _slide_qa(
    prs: Presentation, *, firm_profile: dict[str, Any], bid_slug: str
) -> None:
    person = primary_personnel(firm_profile)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, NAVY)
    _add_rect(
        slide,
        left=Inches(0), top=Inches(0),
        width=SLIDE_WIDTH, height=Inches(0.18),
        fill=GOLD,
    )

    _add_text(
        slide,
        text="Q & A",
        left=Inches(0.5), top=Inches(1.0),
        width=Inches(12.333), height=Inches(1.5),
        font_size=Pt(60), bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        font_name=PPTX_FONT_HEADING,
    )

    _add_text(
        slide,
        text="THANK YOU FOR YOUR CONSIDERATION",
        left=Inches(0.5), top=Inches(2.7),
        width=Inches(12.333), height=Inches(0.5),
        font_size=Pt(13), bold=True, color=GOLD, align=PP_ALIGN.CENTER,
    )

    _add_rect(
        slide,
        left=Inches(6.166), top=Inches(3.4),
        width=Inches(1.0), height=Inches(0.04),
        fill=GOLD,
    )

    contact_lines = [
        person.get("name", ""),
        person.get("title", ""),
        person.get("email", ""),
        person.get("phone", ""),
        firm_profile.get("website", ""),
    ]
    for i, line in enumerate(contact_lines):
        _add_text(
            slide,
            text=line,
            left=Inches(0.5),
            top=Inches(3.7 + i * 0.42),
            width=Inches(12.333), height=Inches(0.4),
            font_size=Pt(15), color=WHITE, align=PP_ALIGN.CENTER,
        )

    _add_text(
        slide,
        text=f"Bid {bid_slug}   ·   {date.today().isoformat()}",
        left=Inches(0.5), top=Inches(7.05),
        width=Inches(12.333), height=Inches(0.3),
        font_size=Pt(9), color=GOLD, align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Slide-builder primitives
# ---------------------------------------------------------------------------


def _body_slide(
    prs: Presentation,
    *,
    title: str,
    bullets: list[tuple[bool, str]] | None = None,
    body_lines: list[str] | None = None,
    bid_slug: str,
    bold_first_phrase: bool = False,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, WHITE)
    _add_title_bar(slide, title)

    body = slide.shapes.add_textbox(
        left=Inches(0.5), top=Inches(1.4),
        width=Inches(12.333), height=Inches(5.4),
    ).text_frame
    body.word_wrap = True

    items = (
        bullets
        if bullets is not None
        else [(False, ln) for ln in (body_lines or [])]
    )
    if not items:
        items = [(False, "(no content)")]

    for i, (is_red, text) in enumerate(items):
        if i == 0:
            para = body.paragraphs[0]
        else:
            para = body.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.level = 0
        runs_added = False
        if bold_first_phrase and ":" in text:
            head, rest = text.split(":", 1)
            r1 = para.add_run()
            r1.text = head + ":"
            _style_run(
                r1,
                font_size=Pt(17), bold=True, color=NAVY,
                font_name=PPTX_FONT_HEADING,
            )
            r2 = para.add_run()
            r2.text = rest
            _style_run(
                r2, font_size=Pt(17), color=NEAR_BLACK, font_name=PPTX_FONT_BODY
            )
            runs_added = True
        else:
            for is_ph, frag in split_placeholders(text):
                if not frag:
                    continue
                r = para.add_run()
                r.text = frag
                _style_run(
                    r,
                    font_size=Pt(17),
                    bold=is_red or is_ph,
                    color=RED if (is_red or is_ph) else NEAR_BLACK,
                    font_name=PPTX_FONT_BODY,
                )
                runs_added = True
        if not runs_added:
            r = para.add_run()
            r.text = text
            _style_run(
                r, font_size=Pt(17), color=NEAR_BLACK, font_name=PPTX_FONT_BODY
            )

    _add_footer(slide, bid_slug=bid_slug, slide_number=len(prs.slides))


def _stat_slide(
    prs: Presentation,
    *,
    title: str,
    stat_label: str,
    stat_value: str,
    bullets: list[tuple[bool, str]],
    bid_slug: str,
) -> None:
    """Slide with a single big stat plus 3 supporting bullets — used for
    the price-proposal summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, WHITE)
    _add_title_bar(slide, title)

    # Stat box (gold accent + light gray bg)
    box_left = Inches(0.5)
    box_top = Inches(1.55)
    box_w = Inches(12.333)
    box_h = Inches(2.0)

    _add_rect(
        slide,
        left=box_left, top=box_top,
        width=box_w, height=box_h,
        fill=GRAY_BG,
    )
    # Gold left rule
    _add_rect(
        slide,
        left=box_left, top=box_top,
        width=Inches(0.08), height=box_h,
        fill=GOLD,
    )

    _add_text(
        slide,
        text=stat_label,
        left=Inches(0.85), top=Inches(1.75),
        width=Inches(11.5), height=Inches(0.4),
        font_size=Pt(11), bold=True, color=SLATE,
    )
    _add_text(
        slide,
        text=stat_value,
        left=Inches(0.85), top=Inches(2.15),
        width=Inches(11.5), height=Inches(1.4),
        font_size=Pt(48), bold=True, color=NAVY,
        font_name=PPTX_FONT_HEADING,
    )

    body_top = Inches(3.85)
    body = slide.shapes.add_textbox(
        left=Inches(0.5), top=body_top,
        width=Inches(12.333), height=Inches(3.0),
    ).text_frame
    body.word_wrap = True
    for i, (is_red, text) in enumerate(bullets):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        for is_ph, frag in split_placeholders(text):
            if not frag:
                continue
            r = para.add_run()
            r.text = frag
            _style_run(
                r,
                font_size=Pt(15),
                bold=is_red or is_ph,
                color=RED if (is_red or is_ph) else NEAR_BLACK,
                font_name=PPTX_FONT_BODY,
            )

    _add_footer(slide, bid_slug=bid_slug, slide_number=len(prs.slides))


def _table_slide(
    prs: Presentation,
    *,
    title: str,
    headers: list[str],
    rows: list[list[str]] | list[tuple[str, ...]],
    bid_slug: str,
    col_widths: list[Emu] | None = None,
    font_size: Pt = Pt(11),
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, WHITE)
    _add_title_bar(slide, title)

    n_cols = len(headers)
    rows_norm: list[list[str]] = []
    for r in rows:
        rl = list(r)
        while len(rl) < n_cols:
            rl.append("")
        rows_norm.append(rl[:n_cols])

    n_rows = len(rows_norm) + 1
    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(12.333)
    height = Inches(5.0)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths and len(col_widths) == n_cols:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = h
        _style_run(
            run,
            font_size=Pt(12),
            bold=True,
            color=WHITE,
            font_name=PPTX_FONT_HEADING,
        )
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r_idx, row in enumerate(rows_norm, start=1):
        is_alt = r_idx % 2 == 0
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = GRAY_BG if is_alt else WHITE
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            para = tf.paragraphs[0]
            for is_ph, frag in split_placeholders(val):
                if not frag:
                    continue
                run = para.add_run()
                run.text = frag
                _style_run(
                    run,
                    font_size=font_size,
                    bold=is_ph,
                    color=RED if is_ph else NEAR_BLACK,
                    font_name=PPTX_FONT_BODY,
                )
            if not para.runs:
                run = para.add_run()
                run.text = ""

    _add_footer(slide, bid_slug=bid_slug, slide_number=len(prs.slides))


def _add_title_bar(slide, title: str) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=Inches(0),
        width=SLIDE_WIDTH, height=Inches(1.0),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = title
    _style_run(
        run,
        font_size=Pt(24),
        bold=True,
        color=WHITE,
        font_name=PPTX_FONT_HEADING,
    )

    # Gold accent stripe under the title bar
    _add_rect(
        slide,
        left=Inches(0), top=Inches(1.0),
        width=SLIDE_WIDTH, height=Inches(0.06),
        fill=GOLD,
    )


def _add_footer(slide, *, bid_slug: str, slide_number: int) -> None:
    box = slide.shapes.add_textbox(
        left=Inches(0.5), top=Inches(7.05),
        width=Inches(8.5), height=Inches(0.35),
    )
    tf = box.text_frame
    tf.margin_left = Inches(0)
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    r1 = para.add_run()
    r1.text = f"Blue Print Constructs   ·   {bid_slug}"
    _style_run(r1, font_size=Pt(9), color=FOOTER_GRAY)

    box_r = slide.shapes.add_textbox(
        left=Inches(10.5), top=Inches(7.05),
        width=Inches(2.333), height=Inches(0.35),
    )
    tf_r = box_r.text_frame
    para_r = tf_r.paragraphs[0]
    para_r.alignment = PP_ALIGN.RIGHT
    r2 = para_r.add_run()
    r2.text = f"Slide {slide_number}"
    _style_run(r2, font_size=Pt(9), color=FOOTER_GRAY)


def _add_text(
    slide,
    *,
    text: str,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    font_size: Pt = Pt(14),
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = NEAR_BLACK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = PPTX_FONT_BODY,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    _style_run(
        run,
        font_size=font_size,
        bold=bold,
        italic=italic,
        color=color,
        font_name=font_name,
    )


def _add_rect(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill: RGBColor,
) -> None:
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.fill.background()


def _style_run(
    run,
    *,
    font_size: Pt = Pt(14),
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = NEAR_BLACK,
    font_name: str = PPTX_FONT_BODY,
) -> None:
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if font_name:
        run.font.name = font_name


def _set_background(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=Inches(0),
        width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _truncate(text: str, max_len: int) -> str:
    text = collapse_whitespace(strip_blockquote_prefix(text or ""))
    if len(text) <= max_len:
        return text
    return text[: max_len - 1].rstrip() + "…"


def _chunk_text_for_slide(body: str, *, max_chars: int) -> list[str]:
    body = strip_blockquote_prefix(body).strip()
    if not body:
        return []
    chunks: list[str] = []
    cursor = 0
    while cursor < len(body):
        end = min(cursor + max_chars, len(body))
        if end < len(body):
            window = body[cursor:end]
            last_break = max(
                window.rfind(". "),
                window.rfind("\n"),
                window.rfind("; "),
            )
            if last_break > max_chars * 0.6:
                end = cursor + last_break + 1
        chunks.append(body[cursor:end].strip())
        cursor = end
    return chunks
